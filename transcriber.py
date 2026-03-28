#!/usr/bin/env python3
"""Transcribe UADE class recordings and generate smart summaries.

Finds .mp4 files in ~/UADE/4to cuatrimestre/*/05_Grabaciones/,
transcribes with mlx-whisper (local, Apple Silicon), and generates
study-oriented summaries via Claude Code CLI with course material context.
Summaries go to 02_Apuntes_Personales/.
"""

import argparse
import os
import re
import sqlite3
import threading
from concurrent.futures import ThreadPoolExecutor, Future
from datetime import datetime
from pathlib import Path

import config

# --- Config ---

BASE_DIR = config.BASE_DIR
DB_PATH = config.DB_PATH

# Materia names for display (derived from folder names)
MATERIA_DISPLAY = {
    "Inteligencia_Artificial_Aplicada": "IA Aplicada",
    "Proceso_de_Desarrollo_de_Software": "PDS",
    "Desarrollo_de_Aplicaciones": "Desarrollo de Aplicaciones I",
    "Ingenieria_de_Datos_II": "Ingeniería de Datos II",
}


# --- DB ---


def init_db():
    conn = sqlite3.connect(str(DB_PATH), check_same_thread=False)
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("""
        CREATE TABLE IF NOT EXISTS transcriptions (
            mp4_path    TEXT PRIMARY KEY,
            txt_path    TEXT NOT NULL,
            summary_path TEXT,
            transcribed_at TEXT NOT NULL,
            summarized_at TEXT,
            context_hash TEXT
        )
    """)
    # Migrar tabla vieja si falta la columna context_hash
    try:
        conn.execute("SELECT context_hash FROM transcriptions LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE transcriptions ADD COLUMN context_hash TEXT")
    conn.commit()
    return conn


def is_transcribed(conn, mp4_path: str) -> bool:
    row = conn.execute(
        "SELECT 1 FROM transcriptions WHERE mp4_path=?", (mp4_path,)
    ).fetchone()
    return row is not None


def needs_resummarize(conn, mp4_path: str, current_hash: str) -> bool:
    """Retorna True si el resumen necesita regenerarse (material cambio)."""
    row = conn.execute(
        "SELECT context_hash, summary_path FROM transcriptions WHERE mp4_path=?",
        (mp4_path,)
    ).fetchone()
    if not row:
        return False
    old_hash, summary_path = row
    if not summary_path:
        return True  # nunca se resumio
    if old_hash != current_hash:
        return True  # material cambio
    return False


def record_transcription(conn, mp4_path: str, txt_path: str,
                         summary_path: str | None = None,
                         context_hash: str | None = None):
    now = datetime.now().isoformat()
    conn.execute(
        """INSERT OR REPLACE INTO transcriptions
           (mp4_path, txt_path, summary_path, transcribed_at, summarized_at, context_hash)
           VALUES (?, ?, ?, ?, ?, ?)""",
        (mp4_path, txt_path, summary_path, now,
         now if summary_path else None, context_hash),
    )
    conn.commit()


# --- Text extraction ---


def extract_text_from_file(path: Path) -> str | None:
    """Extract text from PDF, PPTX, or plain text files."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = [p.extract_text() or "" for p in reader.pages]
            return "\n\n".join(pages).strip()
        elif suffix in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n\n".join(texts).strip()
        elif suffix in (".txt", ".md"):
            return path.read_text(encoding="utf-8").strip()
        elif suffix == ".docx":
            # Basic docx extraction via python-pptx's oxml (no extra dep)
            import zipfile
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(str(path)) as z:
                xml_content = z.read("word/document.xml")
            tree = ET.fromstring(xml_content)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            return "\n".join(
                t.text for t in tree.iter(f"{{{ns['w']}}}t") if t.text
            ).strip()
    except Exception as e:
        log(f"    WARN: no se pudo extraer texto de {path.name}: {e}")
    return None


# --- Helpers ---


def log(msg: str):
    ts = datetime.now().strftime("%H:%M:%S")
    print(f"[{ts}] {msg}")


def find_mp4s() -> list[Path]:
    """Find .mp4 files in 05_Grabaciones/ that haven't been transcribed."""
    results = []
    for materia_dir in sorted(BASE_DIR.iterdir()):
        if not materia_dir.is_dir():
            continue
        grab_dir = materia_dir / config.FOLDERS["grabacion"]
        if not grab_dir.exists():
            continue
        for mp4 in sorted(grab_dir.glob("*.mp4")):
            txt = mp4.with_suffix(".txt")
            if txt.exists():
                continue
            results.append(mp4)
    return results


def extract_class_num(mp4_path: Path) -> int | None:
    """Extract class number from GRAB_XX_ prefix in filename."""
    m = re.match(r'GRAB_(\d+)', mp4_path.stem)
    if m:
        num = int(m.group(1))
        return num if num > 0 else None
    return None


def extract_date(mp4_path: Path) -> str | None:
    """Extract date from GRAB_XX_YYYYMMDD_ in filename."""
    m = re.search(r'GRAB_\d+_(\d{8})', mp4_path.stem)
    if m:
        d = m.group(1)
        return f"{d[:4]}-{d[4:6]}-{d[6:8]}"
    return None


def compute_context_hash(materia_dir: Path, class_num: int | None) -> str:
    """Compute hash of context files (slides + cronograma) for change detection."""
    import hashlib
    h = hashlib.md5()
    material_dir = materia_dir / config.FOLDERS["material"]
    extra_dir = materia_dir / config.FOLDERS["extra"]

    # Hash slides de esta clase
    if class_num and material_dir.exists():
        prefix = f"CLASE_{class_num:02d}_"
        for f in sorted(material_dir.iterdir()):
            if f.name.startswith(prefix) and f.suffix.lower() in (".pdf", ".pptx", ".ppt", ".docx"):
                h.update(f.name.encode())
                h.update(str(f.stat().st_size).encode())
                h.update(str(f.stat().st_mtime).encode())

    # Hash cronograma
    if extra_dir.exists():
        for f in sorted(extra_dir.iterdir()):
            if f.suffix.lower() == ".pdf" and ("3.4." in f.name or "cronograma" in f.name.lower()):
                h.update(f.name.encode())
                h.update(str(f.stat().st_size).encode())
                h.update(str(f.stat().st_mtime).encode())

    return h.hexdigest()


def find_class_context(materia_dir: Path, class_num: int | None) -> dict:
    """Find contextual materials for a specific class.

    Returns dict with keys: slides, cronograma, prev_summary, materia_name
    Each value is extracted text (str) or None.
    """
    ctx = {
        "slides": None,
        "cronograma": None,
        "prev_summary": None,
        "materia_name": MATERIA_DISPLAY.get(materia_dir.name, materia_dir.name),
    }

    material_dir = materia_dir / config.FOLDERS["material"]
    extra_dir = materia_dir / config.FOLDERS["extra"]
    apuntes_dir = materia_dir / config.FOLDERS["apuntes"]

    # 1. Slides for this specific class
    if class_num and material_dir.exists():
        prefix = f"CLASE_{class_num:02d}_"
        slide_texts = []
        for f in sorted(material_dir.iterdir()):
            if f.name.startswith(prefix) and f.suffix.lower() in (".pdf", ".pptx", ".ppt", ".docx"):
                text = extract_text_from_file(f)
                if text:
                    slide_texts.append(f"### {f.name}\n{text}")
        if slide_texts:
            ctx["slides"] = "\n\n".join(slide_texts)

    # 2. Cronograma / plan de materia
    if extra_dir.exists():
        for f in sorted(extra_dir.iterdir()):
            if f.suffix.lower() == ".pdf" and ("3.4." in f.name or "cronograma" in f.name.lower()):
                text = extract_text_from_file(f)
                if text:
                    ctx["cronograma"] = text
                    break

    # 3. Previous class summary for continuity
    if class_num and class_num > 1 and apuntes_dir.exists():
        prev_pattern = f"GRAB_{class_num - 1:02d}_"
        for f in sorted(apuntes_dir.iterdir()):
            if f.name.startswith(prev_pattern) and f.name.endswith("_resumen.md"):
                text = f.read_text(encoding="utf-8").strip()
                if text:
                    # Only keep the headers and first lines to save context
                    lines = text.split("\n")
                    ctx["prev_summary"] = "\n".join(lines[:80])
                    break

    return ctx


def assess_transcription_quality(text: str, mp4_path: Path) -> str | None:
    """Analiza la calidad de la transcripcion. Retorna advertencia o None."""
    issues = []

    size_mb = mp4_path.stat().st_size / (1024 * 1024)
    chars_per_mb = len(text) / max(size_mb, 1)

    # Ratio chars/MB muy bajo = poco contenido para el tamaño del video
    # Una clase normal de 2hs (~400MB) genera ~80K-100K chars (~200+ chars/MB)
    # Un audio malo o con mucho silencio genera mucho menos
    if size_mb > 50 and chars_per_mb < 80:
        issues.append(f"poco contenido transcripto para el tamaño del video ({chars_per_mb:.0f} chars/MB, normal: >150)")

    # Repeticiones excesivas (señal de audio malo / eco / loop)
    words = text.split()
    if len(words) > 50:
        repeats = 0
        for i in range(2, len(words)):
            if words[i] == words[i-1] == words[i-2] and len(words[i]) > 2:
                repeats += 1
        repeat_ratio = repeats / len(words)
        if repeat_ratio > 0.02:
            issues.append("repeticiones excesivas detectadas (posible audio con ruido o eco)")

    if issues:
        return "ADVERTENCIA: Calidad de audio baja — " + "; ".join(issues) + "."
    return None


def transcribe(mp4_path: Path, model: str = "") -> str:
    """Transcribe an mp4 file using auto-detected Whisper backend."""
    size_mb = mp4_path.stat().st_size / (1024 * 1024)
    backend = config.detect_whisper_backend()
    log(f"  Transcribiendo: {mp4_path.name} ({size_mb:.0f} MB) [{backend}]")

    # Heartbeat: print elapsed time every 30s so user knows it's alive
    stop_heartbeat = threading.Event()
    start = datetime.now()

    def heartbeat():
        while not stop_heartbeat.wait(30):
            elapsed = datetime.now() - start
            mins = int(elapsed.total_seconds() // 60)
            secs = int(elapsed.total_seconds() % 60)
            log(f"  ... transcribiendo ({mins}m {secs}s)")

    hb = threading.Thread(target=heartbeat, daemon=True)
    hb.start()

    try:
        text = config.transcribe_audio(str(mp4_path), model=model)
    finally:
        stop_heartbeat.set()
        hb.join(timeout=1)

    elapsed = datetime.now() - start
    mins = int(elapsed.total_seconds() // 60)
    secs = int(elapsed.total_seconds() % 60)
    log(f"  Transcripción completada en {mins}m {secs}s")
    return text


def summarize(transcript: str, mp4_path: Path, class_num: int | None,
              class_date: str | None, ctx: dict,
              quality_warning: str | None = None) -> str:
    """Generate a study-oriented class summary."""

    materia = ctx["materia_name"]
    clase_label = f"Clase {class_num}" if class_num else "Clase"
    fecha_label = class_date or "fecha desconocida"

    # Build context sections
    context_parts = []

    if ctx["cronograma"]:
        context_parts.append(
            f"## Cronograma/Plan de la materia\n{ctx['cronograma'][:3000]}"
        )

    if ctx["slides"]:
        context_parts.append(
            f"## Slides de esta clase\n{ctx['slides'][:8000]}"
        )

    if ctx["prev_summary"]:
        context_parts.append(
            f"## Resumen de la clase anterior\n{ctx['prev_summary']}"
        )

    context_block = "\n\n---\n\n".join(context_parts) if context_parts else ""

    prompt = f"""Sos un asistente académico experto. Generá un resumen estructurado de esta clase
universitaria para estudio. El resumen debe ser útil para: preparar parciales, ponerse al día
si se faltó, armar apuntes propios, y no perderse entregas.

Materia: {materia}
Clase: {clase_label} ({fecha_label})

{f"--- MATERIAL DE CONTEXTO ---{chr(10)}{chr(10)}{context_block}{chr(10)}{chr(10)}" if context_block else ""}--- TRANSCRIPCIÓN DE LA CLASE ---

{transcript}

--- INSTRUCCIONES ---

Generá el resumen con EXACTAMENTE este formato markdown:

# {materia} — {clase_label} ({fecha_label})

## Ubicación en el programa
- Unidad/tema según el cronograma (si se proporcionó)
- Conexión con la clase anterior (si hay resumen previo)

## Temas explicados
- Cada tema con una explicación concisa pero precisa
- Priorizá lo que el profe desarrolló en detalle

## Conceptos clave para el examen
- Concepto: definición precisa tal como la dio el profesor
- Prestá atención especial a frases como "esto es importante", "esto cae", "recuerden que", "esto lo vamos a evaluar"

## Ejemplos y casos prácticos
- Cada ejemplo mencionado y qué concepto ilustra

## Lo que dijo el profe (citas relevantes)
> Frases textuales que enfatizan algo importante o dan pistas sobre evaluaciones

## Correspondencia con slides
- Qué slides/documentos corresponden a cada tema (si se proporcionó material)

## Tareas y entregas
- [ ] Descripción de la tarea 📅 YYYY-MM-DD
- Si no mencionó fecha exacta, poner la fecha estimada según contexto
- Si no hay tareas: "No se asignaron tareas en esta clase."

## Fechas y deadlines mencionados
- Parcial/entrega/evento: fecha
- Si no se mencionaron: "No se mencionaron fechas."

## Dudas para revisar
- Puntos ambiguos, cosas que el profe dijo que iba a retomar, o temas que quedaron incompletos

---

REGLAS:
- Escribí en español argentino
- No inventes información que no esté en la transcripción
- Las citas del profe deben ser lo más textuales posible
- El formato de tareas debe ser compatible con Obsidian Tasks
- Sé completo pero no redundante"""

    summary = config.llm_complete(prompt)
    if quality_warning:
        summary = f"> **{quality_warning}**\n\n{summary}"
    return summary


# --- Main ---


def main():
    parser = argparse.ArgumentParser(description="Transcribe UADE class recordings")
    parser.add_argument("--no-summary", action="store_true",
                        help="Solo transcribir, sin resumen")
    parser.add_argument("--model", type=str, default="",
                        help="Modelo whisper (default: medium)")
    parser.add_argument("--file", type=str,
                        help="Procesar solo este archivo .mp4")
    args = parser.parse_args()

    conn = init_db()
    whisper_model = args.model

    if args.file:
        mp4s = [Path(args.file)]
    else:
        mp4s = find_mp4s()

    if not mp4s:
        log("No hay .mp4 pendientes de transcripción.")
        # Mostrar status actual
        show_pipeline_status()
        return

    log(f"Encontrados {len(mp4s)} videos pendientes")

    # Auto-detectar si hay LLM disponible para resumenes
    skip_summary = args.no_summary
    if not skip_summary and not config.has_llm():
        log("AVISO: No hay LLM configurado. Solo se hara transcripcion (sin resumenes).")
        log("  Para resumenes, configurar: Claude Code, GEMINI_API_KEY, u Ollama.")
        skip_summary = True

    # Thread pool for summarization (runs in parallel with next transcription)
    executor = ThreadPoolExecutor(max_workers=1) if not skip_summary else None
    pending_summary: Future | None = None

    def wait_pending_summary():
        nonlocal pending_summary
        if pending_summary is None:
            return
        try:
            pending_summary.result()
        except Exception as e:
            log(f"  ERROR en resumen pendiente: {e}")
        pending_summary = None

    def submit_summary(text, mp4, mp4_path_str, txt_path, summary_path, materia_dir,
                       is_regen=False):
        nonlocal pending_summary
        wait_pending_summary()

        class_num = extract_class_num(mp4)
        class_date = extract_date(mp4)

        def _do_summary():
            label = "Regenerando" if is_regen else "Recopilando contexto para"
            log(f"  {label} {mp4.name}...")
            ctx = find_class_context(materia_dir, class_num)
            ctx_hash = compute_context_hash(materia_dir, class_num)
            parts = [k for k in ("slides", "cronograma", "prev_summary") if ctx[k]]
            log(f"  Contexto: {', '.join(parts) if parts else 'solo transcripción'}")
            quality_warn = assess_transcription_quality(text, mp4)
            if quality_warn:
                log(f"  {quality_warn}")
            log(f"  Generando resumen con {config.detect_llm_provider()}...")
            summary = summarize(text, mp4, class_num, class_date, ctx,
                                quality_warning=quality_warn)
            summary_path.write_text(summary, encoding="utf-8")
            log(f"  OK: {summary_path.name}")
            record_transcription(conn, mp4_path_str, str(txt_path),
                                 str(summary_path), context_hash=ctx_hash)

        pending_summary = executor.submit(_do_summary)

    for mp4 in mp4s:
        if is_transcribed(conn, str(mp4)):
            log(f"  SKIP (ya en DB): {mp4.name}")
            continue

        txt_path = mp4.with_suffix(".txt")
        materia_dir = mp4.parent.parent
        apuntes_dir = materia_dir / config.FOLDERS["apuntes"]
        apuntes_dir.mkdir(parents=True, exist_ok=True)
        summary_path = apuntes_dir / (mp4.stem + "_resumen.md")

        # Transcribe (GPU)
        try:
            text = transcribe(mp4, model=whisper_model)
        except Exception as e:
            log(f"  ERROR transcribiendo {mp4.name}: {e}")
            continue

        txt_path.write_text(text, encoding="utf-8")
        log(f"  OK: {txt_path.name} ({len(text)} chars)")

        # Register transcription immediately (summary updates later)
        record_transcription(conn, str(mp4), str(txt_path))

        # Summarize (network — runs in parallel with next transcription)
        if not skip_summary and executor:
            submit_summary(text, mp4, str(mp4), txt_path, summary_path, materia_dir)

    # Wait for last summary
    wait_pending_summary()

    # Check for stale summaries (material changed since last summary)
    if not skip_summary and executor:
        stale_count = 0
        for materia_dir in sorted(BASE_DIR.iterdir()):
            if not materia_dir.is_dir():
                continue
            grab_dir = materia_dir / config.FOLDERS["grabacion"]
            if not grab_dir.exists():
                continue
            for mp4 in sorted(grab_dir.glob("*.mp4")):
                txt_path = mp4.with_suffix(".txt")
                if not txt_path.exists():
                    continue
                class_num = extract_class_num(mp4)
                ctx_hash = compute_context_hash(materia_dir, class_num)
                if needs_resummarize(conn, str(mp4), ctx_hash):
                    stale_count += 1
                    apuntes_dir = materia_dir / config.FOLDERS["apuntes"]
                    apuntes_dir.mkdir(parents=True, exist_ok=True)
                    summary_path = apuntes_dir / (mp4.stem + "_resumen.md")
                    text = txt_path.read_text(encoding="utf-8")
                    log(f"  Material nuevo detectado para {mp4.name}")
                    submit_summary(text, mp4, str(mp4), txt_path, summary_path,
                                   materia_dir, is_regen=True)
        if stale_count:
            wait_pending_summary()
            log(f"  {stale_count} resumenes regenerados por cambio de material")

    if executor:
        executor.shutdown(wait=True)

    conn.close()

    # Consolidar tareas de todos los resumenes
    if not skip_summary:
        consolidate_tasks()

    show_pipeline_status()
    log("Listo.")


def show_pipeline_status():
    """Muestra un resumen del estado del pipeline."""
    total_mp4 = 0
    total_txt = 0
    total_resumen = 0
    for materia_dir in sorted(BASE_DIR.iterdir()):
        if not materia_dir.is_dir():
            continue
        grab_dir = materia_dir / config.FOLDERS["grabacion"]
        apuntes_dir = materia_dir / config.FOLDERS["apuntes"]
        if grab_dir.exists():
            mp4s = list(grab_dir.glob("*.mp4"))
            txts = list(grab_dir.glob("*.txt"))
            total_mp4 += len(mp4s)
            total_txt += len(txts)
        if apuntes_dir.exists():
            resums = list(apuntes_dir.glob("*_resumen.md"))
            total_resumen += len(resums)
    log(f"  Estado: {total_mp4} grabaciones, {total_txt} transcripciones, {total_resumen} resumenes")


def consolidate_tasks():
    """Parsea todos los _resumen.md y genera un tareas.md por materia."""
    import re as _re
    task_pattern = _re.compile(r'^- \[ \] .+', _re.MULTILINE)

    for materia_dir in sorted(BASE_DIR.iterdir()):
        if not materia_dir.is_dir():
            continue
        apuntes_dir = materia_dir / config.FOLDERS["apuntes"]
        if not apuntes_dir.exists():
            continue

        all_tasks = []
        for resumen in sorted(apuntes_dir.glob("*_resumen.md")):
            text = resumen.read_text(encoding="utf-8")
            tasks = task_pattern.findall(text)
            if tasks:
                clase_name = resumen.stem.replace("_resumen", "")
                all_tasks.append(f"### {clase_name}\n")
                all_tasks.extend(tasks)
                all_tasks.append("")

        if all_tasks:
            tareas_path = materia_dir / "tareas.md"
            content = f"# Tareas — {materia_dir.name}\n\n"
            content += "Generado automaticamente por el pipeline.\n\n"
            content += "\n".join(all_tasks) + "\n"
            tareas_path.write_text(content, encoding="utf-8")
            log(f"  Tareas consolidadas: {tareas_path.name} ({len([t for t in all_tasks if t.startswith('- [ ]')])} tareas)")


if __name__ == "__main__":
    main()
